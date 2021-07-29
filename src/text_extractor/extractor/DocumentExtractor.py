from ..constants import (
    EXTRACTOR_MODE_DOC,
)
from .Extractor import Extractor
import PyPDF2
import docx
import pptx
import os
import re

class DocumentExtractor(Extractor):
    def __init__(self):
        self.mode = EXTRACTOR_MODE_DOC

    def read_input(self, filepath):
        output = ""
        ext = os.path.splitext(filepath)[1]
        if ext == ".pdf":
            output = self.read_input_pdf(filepath)
            print(output)
        elif ext == ".docx":
            output = self.read_input_docx(filepath)
            print(output)
        elif ext == ".pptx":
            output = self.read_input_pptx(filepath)
            print(output)
        else:
            return False
        return output
    
    def read_input_pdf(self, filepath):
        output = ""
        with open(filepath, 'rb') as stream:
            pdfReader = PyPDF2.PdfFileReader(stream)
            for page in range(pdfReader.numPages):
                pageObj = pdfReader.getPage(page)
                output = output + " " + pageObj.extractText()
        return self.process_raw_text(output)

    def read_input_pptx(self, filepath):
        presentation = pptx.Presentation(filepath)
        fullText = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                try:
                    fullText.append(shape.text)
                except:
                    continue
        return self.process_raw_text('\n'.join(fullText))

    def read_input_docx(self, filepath):
        doc = docx.Document(filepath)
        fullText = []
        for paragraph in doc.paragraphs:
            fullText.append(paragraph.text)
        # for table in doc.tables:
        #     for row in table.rows:
        #         for cell in row.cells:
        #             fullText.append(cell.text)
        return self.process_raw_text('\n'.join(fullText))
        
    def process_raw_text(self, text):
        text = text.replace("\n", "")
        text = re.sub(r'[\\s]{2,}', '\n', text)
        text = re.sub(r'([a-z]{1})([A-Z]{1})', r'\1 \2', text)
        return text
    

        

        


